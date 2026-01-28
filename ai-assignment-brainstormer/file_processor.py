# File processing utilities for extracting text from various file types
import io
from PIL import Image
import pytesseract
from PyPDF2 import PdfReader
from pptx import Presentation


def extract_text_from_pdf(file) -> str:
    """Extract text from a PDF file."""
    try:
        reader = PdfReader(file)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts) if text_parts else "No text found in PDF."
    except Exception as e:
        return f"Error reading PDF: {str(e)}"


def extract_text_from_pptx(file) -> str:
    """Extract text from a PowerPoint file."""
    try:
        prs = Presentation(file)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts) if text_parts else "No text found in PowerPoint."
    except Exception as e:
        return f"Error reading PowerPoint: {str(e)}"


def extract_text_from_image(file) -> str:
    """Extract text from an image using OCR (supports handwritten text too)."""
    try:
        image = Image.open(file)
        # Use pytesseract for OCR
        text = pytesseract.image_to_string(image, lang='eng+ind')  # English + Indonesian
        return text.strip() if text.strip() else "No text detected in image."
    except Exception as e:
        return f"Error reading image: {str(e)}"


def process_uploaded_file(uploaded_file) -> tuple[str, str]:
    """
    Process an uploaded file and extract its text content.
    
    Returns:
        tuple: (extracted_text, file_type)
    """
    if uploaded_file is None:
        return "", ""
    
    file_name = uploaded_file.name.lower()
    
    # Determine file type and extract text
    if file_name.endswith('.pdf'):
        text = extract_text_from_pdf(uploaded_file)
        file_type = "PDF"
    elif file_name.endswith(('.pptx', '.ppt')):
        text = extract_text_from_pptx(uploaded_file)
        file_type = "PowerPoint"
    elif file_name.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp')):
        text = extract_text_from_image(uploaded_file)
        file_type = "Image"
    else:
        text = "Unsupported file type. Please upload PDF, PPTX, or image files."
        file_type = "Unknown"
    
    return text, file_type
